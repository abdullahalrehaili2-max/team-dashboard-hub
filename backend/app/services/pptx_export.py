"""PPTX export using python-pptx."""

import os
import tempfile
from sqlalchemy.orm import Session
from app import models


def export_dashboard_pptx(dashboard: models.Dashboard, db: Session) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    pages = db.query(models.Page).filter_by(dashboard_id=dashboard.id).order_by(models.Page.order).all()
    if not pages:
        # Create one slide with dashboard name
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
        tf = txBox.text_frame
        tf.text = dashboard.name
        p = tf.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
    else:
        for page in pages:
            slide = prs.slides.add_slide(blank_layout)
            
            # Page title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
            tf = title_box.text_frame
            tf.text = page.name
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = True

            # Sections as cards
            sections = db.query(models.Section).filter_by(page_id=page.id).order_by(models.Section.order).all()
            col = 0
            row = 0
            for sec in sections[:8]:
                left = Inches(0.5 + col * 3.1)
                top = Inches(1.0 + row * 1.8)
                width = Inches(2.8)
                height = Inches(1.5)

                # Card box
                shape = slide.shapes.add_textbox(left, top, width, height)
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sec.title or sec.indicator_type
                p.font.size = Pt(14)
                p.font.bold = True

                col += 1
                if col >= 4:
                    col = 0
                    row += 1

    out_path = tempfile.mktemp(suffix=".pptx")
    prs.save(out_path)
    return out_path
